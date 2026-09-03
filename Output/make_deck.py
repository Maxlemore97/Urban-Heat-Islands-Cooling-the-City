#!/usr/bin/env python3
"""Discovering & Abstracting deck — Challenge 1, Urban Heat Islands."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
BG      = RGBColor(0xEF, 0xEE, 0xE9)   # concrete paper
SURF    = RGBColor(0xF7, 0xF6, 0xF2)
INK     = RGBColor(0x1A, 0x1B, 0x19)
INK2    = RGBColor(0x4A, 0x4B, 0x45)
INK3    = RGBColor(0x86, 0x86, 0x7C)
RULE    = RGBColor(0xD2, 0xD1, 0xC9)
GOLD    = RGBColor(0xA8, 0x87, 0x1A)   # prevention / day
BLUE    = RGBColor(0x15, 0x70, 0x9F)   # removal / night
VERM    = RGBColor(0xAE, 0x3B, 0x2E)   # self-regulation, F9
GREY    = RGBColor(0x6B, 0x6B, 0x62)   # supporting functions

DISPLAY = "Arial"
BODY    = "Georgia"

def tint(accent, amount=0.13, base=BG):
    return RGBColor(*[int(round(b * (1 - amount) + a * amount))
                      for a, b in zip(accent, base)])

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]

L, R, T = Inches(0.95), Inches(0.75), Inches(0.5)
CW = prs.slide_width - L - R          # content width


# ---------------------------------------------------------------- helpers
def rect(slide, l, t, w, h, fill=None, line=None, lw=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.text_frame.text = ""
    return sh


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def para(tf, text, size=12, bold=False, color=INK, font=BODY, space_before=0,
         space_after=4, spacing=1.15, caps_track=None, align=PP_ALIGN.LEFT,
         italic=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = spacing
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    f = run.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    if caps_track:
        run._r.get_or_add_rPr().set('spc', str(int(caps_track * 100)))
    return p


def base(title_bg=BG):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=title_bg)
    return s


def footer(s, left_text, page=None):
    tf = box(s, L, Inches(6.86), CW - Inches(0.8), Inches(0.3))
    para(tf, left_text, size=8.5, color=INK3, font=DISPLAY, caps_track=1.1,
         first=True, space_after=0)
    rect(s, L, Inches(6.78), CW, Inches(0.01), fill=RULE)
    if page is not None:
        tfp = box(s, prs.slide_width - R - Inches(0.8), Inches(6.86),
                  Inches(0.8), Inches(0.3))
        para(tfp, f"{page:02d}", size=8.5, color=INK3, font=DISPLAY,
             align=PP_ALIGN.RIGHT, first=True, space_after=0)


def eyebrow(s, text, accent):
    rect(s, Inches(0.5), T, Inches(0.055), Inches(6.15), fill=accent)
    tf = box(s, L, T + Inches(0.02), CW, Inches(0.3))
    para(tf, text, size=9, color=accent, font=DISPLAY, bold=True,
         caps_track=1.6, first=True, space_after=0)


# ---------------------------------------------------------------- slides
def title_slide():
    s = base()
    rect(s, L, Inches(2.05), Inches(4.6), Inches(0.045), fill=INK)
    tf = box(s, L, Inches(0.9), CW, Inches(1.1))
    para(tf, "ZHAW · BLOCKWEEK BIONICS 2026 · CHALLENGE 1", size=10,
         color=INK3, font=DISPLAY, bold=True, caps_track=2.0, first=True,
         space_after=14)
    tf2 = box(s, L, Inches(2.35), Inches(9.6), Inches(2.0))
    para(tf2, "Discovering &", size=46, bold=True, color=INK, font=DISPLAY,
         spacing=0.95, first=True, space_after=0)
    para(tf2, "Abstracting", size=46, bold=True, color=INK, font=DISPLAY,
         spacing=0.95, space_after=16)
    tf3 = box(s, L, Inches(4.35), Inches(8.6), Inches(1.4))
    para(tf3, "Twelve biological models for cooling a city, the mechanism "
              "behind each one, and what it would mean for Zurich.",
         size=15, color=INK2, font=BODY, spacing=1.3, first=True)
    cols = [("MODELS FOUND", "12"), ("DE-SELECTED", "4"),
            ("FUNCTIONS COVERED", "F1–F11"), ("REFERENCE CITY", "Zurich")]
    x = L
    for lab, val in cols:
        tfc = box(s, x, Inches(5.75), Inches(2.6), Inches(0.9))
        para(tfc, lab, size=8.5, color=INK3, font=DISPLAY, bold=True,
             caps_track=1.4, first=True, space_after=4)
        para(tfc, val, size=17, bold=True, color=INK, font=DISPLAY,
             space_after=0)
        x += Inches(2.9)
    footer(s, "Cooling the City — urban heat islands")


def statement_slide(page, eyeb, headline, accent, bullets, note=None):
    s = base()
    eyebrow(s, eyeb, accent)
    tf = box(s, L, T + Inches(0.42), Inches(11.2), Inches(1.0))
    para(tf, headline, size=30, bold=True, color=INK, font=DISPLAY,
         spacing=1.02, first=True)
    y = T + Inches(1.35)
    for head, text in bullets:
        tfb = box(s, L, y, Inches(11.4), Inches(0.9))
        para(tfb, head, size=9, color=accent, font=DISPLAY, bold=True,
             caps_track=1.5, first=True, space_after=5)
        para(tfb, text, size=13, color=INK2, font=BODY, spacing=1.25,
             space_after=0)
        y += Inches(1.25)
    if note:
        tfn = box(s, L, Inches(6.15), Inches(11.4), Inches(0.6))
        para(tfn, note, size=11.5, color=INK3, font=BODY, italic=True,
             spacing=1.2, first=True)
    footer(s, eyeb, page)


def model_slide(page, num, organism, common, family, accent, mechanism,
                abstraction, city, link):
    s = base()
    eyebrow(s, f"MODEL {num:02d}   ·   {family}", accent)

    tf = box(s, L, T + Inches(0.40), Inches(11.2), Inches(0.75))
    para(tf, organism, size=29, bold=True, color=INK, font=DISPLAY,
         spacing=1.0, first=True, space_after=3)
    tfs = box(s, L, T + Inches(1.02), Inches(11.2), Inches(0.4))
    para(tfs, common, size=12.5, color=INK3, font=BODY, italic=True,
         first=True, space_after=0)

    # abstraction band — the deliverable of this phase
    band_t = T + Inches(1.62)
    rect(s, L, band_t, CW, Inches(1.20), fill=tint(accent))
    rect(s, L, band_t, Inches(0.05), Inches(1.20), fill=accent)
    tfa = box(s, L + Inches(0.30), band_t + Inches(0.15), CW - Inches(0.6),
              Inches(0.8))
    para(tfa, "ABSTRACTION", size=8.5, color=accent, font=DISPLAY, bold=True,
         caps_track=1.5, first=True, space_after=5)
    para(tfa, abstraction, size=14.5, color=INK, font=BODY, italic=True,
         spacing=1.18, space_after=0)

    # two columns
    col_t = band_t + Inches(1.52)
    cw = (CW - Inches(0.7)) / 2
    for i, (head, text) in enumerate([("THE MECHANISM", mechanism),
                                      ("WHAT IT MEANS FOR THE CITY", city)]):
        x = L + i * (cw + Inches(0.7))
        tfc = box(s, x, col_t, cw, Inches(2.4))
        para(tfc, head, size=8.5, color=INK3, font=DISPLAY, bold=True,
             caps_track=1.4, first=True, space_after=7)
        para(tfc, text, size=12, color=INK2, font=BODY, spacing=1.26,
             space_after=0)

    footer(s, link, page)


# ---------------------------------------------------------------- content
title_slide()

statement_slide(
    2, "THE PHASE", "What this step has to deliver", VERM,
    [("THE SPIRAL",
      "Discovering has two halves. DISCOVER: find biological models and the "
      "evidence for them. ABSTRACT: turn each mechanism into a transferable "
      "design principle — without copying the organism's shape."),
     ("REQUIRED OUTPUT",
      "At least three biological models, an explicit research path showing "
      "how they were found and why others were dropped, mechanisms explained "
      "in depth, and a clean separation of biological observation, abstracted "
      "principle and possible technical interpretation."),
     ("WHAT WE HAVE",
      "Twelve models, each mapped to one of our functions F1–F11, plus four "
      "documented de-selections. The abstraction of every model is stated on "
      "its own slide, in the highlighted band.")],
    note="\"Do not stop at the organism. The mechanism is more important than "
         "the name.\" — Discovering lecture, slide 34")

statement_slide(
    3, "RESEARCH PATH", "How we searched", BLUE,
    [("BY FUNCTION, NOT BY ORGANISM",
      "Every search started from a verb in our function list, never from an "
      "animal. \"How does nature emit heat to the night sky?\" — not "
      "\"how does a termite mound work?\""),
     ("THE TAXONOMY AS THE MAP",
      "F1 reflect → Modify physical state: light/colour. F2 emit → Transform "
      "energy: thermal. F6–F8 water → Capture, Store and Distribute liquids. "
      "F9 is the exception: it is the pairing of Sense signals: temperature "
      "with Process signals: respond to signals."),
     ("FOUR SOURCE TYPES",
      "AskNature strategies and innovations · the Genius of Biome report for "
      "the temperate broadleaf forest, which is Zurich's own biome · primary "
      "papers behind each claim · commercial products where the transfer has "
      "already been made.")],
    note="Discovery is not evidence. AskNature and AI find candidates; the "
         "paper or the measurement is what we cite.")

statement_slide(
    4, "BIOLOGIZED QUESTIONS", "From our functions to questions for nature",
    GOLD,
    [("PREVENTION — F1, F5",
      "How does nature reflect intense sunlight without pigments that bleach "
      "and without a mirror finish? How does nature shade itself with its own "
      "geometry?"),
     ("REMOVAL — F2, F3, F4",
      "How does nature release heat to the night sky? How does nature hand "
      "heat to moving air? How does nature buffer a daily extreme instead of "
      "resisting it?"),
     ("WATER AND REGULATION — F6–F9",
      "How does nature capture, store and meter out scarce water so a reserve "
      "lasts a dry period? And the central one: how does nature make a "
      "response stronger when the load is stronger, without a controller?")],
    note="The challenge brief asks a fourth question we should answer "
         "explicitly: what can we learn from desert organisms about passive "
         "cooling?")

M = [
 dict(num=1, organism="African elephant — ears",
      common="Loxodonta africana · dissipates body heat in hot, exposed habitat",
      family="F2 · F3   REMOVAL", accent=BLUE,
      mechanism="Warm blood is carried from the body core into large, thin "
                "ears. The big, well-supplied surface lets heat escape to the "
                "air, and flapping raises the airflow across it.",
      abstraction="Transfer heat to a large surface and increase the airflow "
                  "around it.",
      city="It separates where heat is picked up from where it is released. A "
           "canyon wall stores heat but cannot get rid of it; a dedicated, "
           "well-exposed surface can.",
      link="asknature.org/strategy/large-ears-aid-cooling/"),
 dict(num=2, organism="Elephant skin — sparse hairs",
      common="Micro-finned surface · heat transfer at low wind speed",
      family="F3   REMOVAL", accent=BLUE,
      mechanism="A thin layer of still air clings to any warm surface and "
                "insulates it. Sparse hairs reach through that boundary layer "
                "so moving air actually reaches the skin. Reported additional "
                "heat loss of up to 23 %.",
      abstraction="Add small protrusions to break the boundary layer and hand "
                  "heat over to passing air.",
      city="Street canyons have low wind speeds, so the insulating boundary "
           "layer is thick. Texture at millimetre scale on façades and "
           "pavements adds no material and no maintenance.",
      link="sites.nd.edu — biomechanics in the wild · asknature.org/strategy/"
           "small-leaves-buffer-insect-eggs-from-heat/"),
 dict(num=3, organism="Tree bark",
      common="e.g. paperbark maple, Acer griseum · keeps the surface cool in sun",
      family="F1 · F2   PREVENTION + REMOVAL", accent=GOLD,
      mechanism="Bark reflects part of the incoming sunlight and emits thermal "
                "radiation efficiently. Rough, layered bark also shades "
                "itself, traps insulating air and improves convection.",
      abstraction="Reflect the incoming energy, emit the stored energy, and "
                  "use texture for shade and airflow.",
      city="One surface doing prevention and removal at once — and bark is the "
           "closest natural analogue to what we are designing: a dead outer "
           "layer protecting a living structure behind it.",
      link="asknature.org/strategy/bark-keeps-surface-cool-under-the-sun/"),
 dict(num=4, organism="Cactus ribs",
      common="Ribbed and barrel cacti · less overheating, less water loss",
      family="F5 · F3 · F2   PREVENTION + REMOVAL", accent=GOLD,
      mechanism="Peaks shade the troughs. The cooler air collecting in the "
                "troughs absorbs heat from the body, rises and is carried "
                "away, while the folds enlarge the radiating area and disturb "
                "the airflow.",
      abstraction="Use folded or ribbed geometry for self-shading, more "
                  "surface and better airflow.",
      city="Costs no material, no water and no chemistry, and there is nothing "
           "to bleach or run out. The conflict is visible change on a "
           "protected façade — the feature size is the negotiable part.",
      link="asknature.org/strategy/shape-shades-and-enhances-heat-radiation/"),
 dict(num=5, organism="Jackrabbit ears · elephant hot spots",
      common="Cooling by varying the transport, not the surface",
      family="F9 · F2   SELF-REGULATION", accent=VERM,
      mechanism="The jackrabbit widens or narrows the vessels in its thin "
                "ears; at around 30 °C air temperature the ears alone carry "
                "all its excess heat, and it costs no water. Elephants open "
                "small cooling windows scattered over the skin, as many as "
                "the heat requires.",
      abstraction="Regulate cooling by varying the transport to a radiating "
                  "surface, and place that radiator only where the conditions "
                  "are best.",
      city="A wall at eye level barely sees the sky — that blocked view is the "
           "physical cause of the night-time heat island. A parapet two floors "
           "up sees it well. A loop driven by density differences needs no "
           "pump and strengthens as the temperature difference grows.",
      link="asknature.org/strategy/how-blood-flow-keeps-jackrabbits-cool/ · "
           "asknature.org/strategy/skin-fine-tunes-internal-temperature/"),
 dict(num=6, organism="Desert snail shell",
      common="Sphincterochila boissieri · survives ground heat up to 50 °C",
      family="F1 · F4   PREVENTION + STORAGE", accent=GOLD,
      mechanism="The shell reflects about 90 % of visible light and about "
                "95 % in the near infrared. When conditions get extreme the "
                "snail seals the opening with a calcified plate and withdraws, "
                "leaving a layer of air between itself and the hot ground.",
      abstraction="Reflect first — then decouple whatever still gets through, "
                  "with a layer of air.",
      city="The 7 K a city is warmer at night is heat released from mass "
           "charged during the day. A reflective skin over a ventilated cavity "
           "stops the masonry charging in the first place — and that is a "
           "standard façade system, not an exotic one.",
      link="asknature.org/strategy/shell-protects-from-heat/"),
 dict(num=7, organism="Morpho butterfly wing",
      common="Brilliantly blue, with no blue pigment at all",
      family="F1   PREVENTION", accent=GOLD,
      mechanism="Transparent layers of chitin and air, a few hundred "
                "nanometres apart, diffract the light and recombine it so that "
                "some wavelengths cancel and others are reinforced. The colour "
                "is in the geometry, and because the structure is disordered "
                "it stays diffuse rather than mirror-like.",
      abstraction="Separate what a surface looks like from what it does with "
                  "heat.",
      city="About half the sun's energy arrives as invisible near infrared. A "
           "façade can keep the ochre of the old town and still reject the "
           "heat — which is what decides whether the measure is ever "
           "permitted, and it answers the glare objection too.",
      link="asknature.org/strategy/wing-scales-cause-light-to-diffract-and-interfere/"),
 dict(num=8, organism="Leaf stomata",
      common="Cooling regulated by the opening itself",
      family="F9 · F8   SELF-REGULATION", accent=VERM,
      mechanism="Two guard cells swell and shrink with their water content and "
                "so open and close the pore. The water is pulled up from the "
                "roots by the evaporation itself — there is no pump anywhere. "
                "The leaf cools itself by more than 10 K.",
      abstraction="Build the control into the opening, so the escape route "
                  "widens as the load rises.",
      city="Evaporation is the only mechanism here that cools the air people "
           "actually breathe, rather than only the surface they walk past. The "
           "binding limit is water: during a heat wave it does not rain.",
      link="Genius of Biome, pp. 62–65 — Day 2 course material"),
 dict(num=9, organism="Camel fur and sweat glands",
      common="Built at MIT as a hydrogel under an aerogel",
      family="F7 · F8   WATER PATH", accent=BLUE,
      mechanism="A hydrogel layer evaporates like a sweat gland, underneath an "
                "aerogel layer that blocks incoming heat but lets the vapour "
                "through, like fur. The pair held a sample 7 K below ambient "
                "for 200 hours — five times longer than the hydrogel alone.",
      abstraction="Cut the heat load first, then spend a little water on what "
                  "is left — and insulate the wet layer so the reserve lasts.",
      city="Zurich has ample rain over a year and almost none during a heat "
           "wave, so our problem is storage and metering, not collection. It "
           "also rules out fog harvesting: fog here is a winter phenomenon.",
      link="asknature.org/innovation/passive-cooling-system-inspired-by-camels/"),
 dict(num=10, organism="Pine cone scales",
      common="Dead tissue, moved entirely by the weather",
      family="F5 · F9   PREVENTION + SELF-REGULATION", accent=VERM,
      mechanism="Each scale is two bonded layers that swell differently with "
                "moisture, so the scale bends. Already built twice: HygroSkin, "
                "whose wooden apertures are shut at 30 % humidity and open at "
                "75 %, and thermobimetal shutters that curl in the sun — "
                "explicitly with no energy and no controls.",
      abstraction="Make moving parts from two layers that respond "
                  "differently, and let the environment do the actuating.",
      city="A shade takes the sun off the person, not just off the wall. And "
           "by retracting when it cools it leaves the night sky open, which "
           "fixed shading does not. The risk is wind: a thermal actuator "
           "cannot feel a storm.",
      link="asknature.org/strategy/pine-cones-open-and-close-in-response-to-weather/"),
 dict(num=11, organism="Honeybee fanning thresholds",
      common="A colony holds 36 °C with no central control",
      family="F9   SELF-REGULATION", accent=VERM,
      mechanism="Every worker starts fanning at a slightly different "
                "temperature, and the spread is genetic — one queen, many "
                "fathers, therefore many thresholds in one nest. As it warms, "
                "a few bees fan, then more, then many.",
      abstraction="Spread many binary elements across a range of trigger "
                  "points, and the sum responds continuously.",
      city="This answers the standard objection to switching materials, that "
           "they are all-or-nothing. Build whatever we build as many small "
           "units with different trigger temperatures, and the response grows "
           "through the afternoon instead of snapping on.",
      link="asknature.org/strategy/varying-response-thresholds-aid-hive-thermoregulation/"),
 dict(num=12, organism="Lichen surface",
      common="Lecanora conizaeoides · waterproof and breathable at once",
      family="F10 · F11   SUPPORTING", accent=GREY,
      mechanism="Roughness at several size scales keeps droplets perched on "
                "the peaks so they roll off, while the channels between them "
                "are coated to stay dry and still let air through to the algae "
                "inside. Dirt and biofilm find almost nothing to hold on to.",
      abstraction="Use geometry rather than chemistry to shed water and stay "
                  "clean.",
      city="The city names soiling as the specific weakness of bright "
           "envelopes. A lotus-effect façade paint has done this since 1999 "
           "without biocide — which matters twice over if we collect that "
           "run-off and evaporate it at head height.",
      link="Genius of Biome, pp. 41–44 · asknature.org/innovation/"
           "paint-inspired-by-lotus-leaves-creates-self-cleaning-and-antifouling-surfaces/"),
]

page = 5
for m in M:
    model_slide(page, m["num"], m["organism"], m["common"], m["family"],
                m["accent"], m["mechanism"], m["abstraction"], m["city"],
                m["link"])
    page += 1

statement_slide(
    17, "DE-SELECTION", "What we looked at and left out", GREY,
    [("THE TERMITE MOUND",
      "The Eastgate building that made it famous uses fans, and the chimney "
      "explanation is contested in the literature (Turner 2008). It is not the "
      "passive system it is sold as."),
     ("FOG HARVESTING — THE NAMIB BEETLE",
      "The most-cited water model in biomimicry, and wrong for our climate: "
      "Zurich's fog is a winter phenomenon, and during a heat wave there is "
      "neither rain nor fog."),
     ("GREEN FAÇADES · GAZELLE CAROTID RETE · TENREC DORMANCY",
      "Greening is already the city's own measures HA 09 and HA 10, and it "
      "needs the most water exactly when we have the least. The other two are "
      "internal physiology and behaviour — a wall has neither.")],
    note="The criterion throughout: it has to be passive, survive frost and UV "
         "for twenty years, and live in a surface.")

statement_slide(
    18, "SYNTHESIS", "What the twelve models have in common", VERM,
    [("ONE PRINCIPLE, FOUND THREE TIMES",
      "Bees vary their thresholds, elephants vary how many cooling windows are "
      "open, a deployable shade varies how much area is spread. All three "
      "regulate how MANY elements act — not how strongly each one acts. That "
      "is our answer to F9."),
     ("PREVENTION IS THE SMALLER HALF",
      "Models 3, 4, 6 and 7 stop heat arriving. Only 1, 2, 5, 8, 9 and 10 get "
      "heat out again — and the urban heat island is 1–2 K by day against up "
      "to 7 K at night. The night belongs to the removal side."),
     ("STRUCTURES TRANSFER BETTER THAN PROCESSES",
      "Most cooling in nature is metabolism or behaviour, and a façade has "
      "neither. The two models that transfer most directly — the snail shell "
      "and the cactus rib — are the two that are structures rather than "
      "processes.")])

statement_slide(
    19, "NEXT", "From principles to concepts", BLUE,
    [("THE MORPHOLOGICAL MATRIX",
      "Rows are our functions, columns are the abstracted principles on these "
      "slides. Reflect: structural colour · ribbed geometry · reflective skin "
      "with an air gap. Remove: passive loop to a roof radiator · boundary-"
      "layer texture · metered evaporation."),
     ("THEN TWO TECHNICALLY DISTINCT CONCEPTS",
      "Not two cosmetic variants. Each with an annotated drawing, an "
      "Engineering Canvas, and two or three critical questions of the form "
      "\"what must be true for this to work?\""),
     ("THE THREE QUESTIONS WE ALREADY KNOW WE HAVE TO ANSWER",
      "Can a passive loop survive a Swiss winter? Can a shading element be "
      "furled by wind without a controller? And can a switching material be "
      "made with a deliberate spread of trigger temperatures?")])

out = ("/private/tmp/claude-501/-Users-maximilian-Library-CloudStorage-"
       "ProtonDrive-Maxlemore97-proton-me-folder-Studium-Informatik-Blockweek-"
       "Bionics-2026/c20395e6-a9ec-4073-8cb8-fcb214bd5b77/scratchpad/"
       "Discovering_and_Abstracting.pptx")
prs.save(out)
print("saved:", out)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
